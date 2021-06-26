from docx import Document
import extract_msg
from pptx import Presentation
from openpyxl import load_workbook
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from io import StringIO

def get_text_pdf(path):
    
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, laparams=laparams)
    fp = open(path, 'rb')
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    password = ""
    maxpages = 0
    pagenos=set()

    for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages,
                                  caching=True, check_extractable=True):
        interpreter.process_page(page)

    text = retstr.getvalue()

    fp.close()
    device.close()
    retstr.close()
    return text

def get_text_docx(path):
    full_text = ''
    document = Document(path)
    for para in document.paragraphs:
        full_text += para.text + '. '
        #Add tables
#     for table in document.tables:
#         for i, row in enumerate(table.rows):
#             for cell in row.cells:
#                 full_text += cell.text
    return full_text

def get_text_pptx(path):
    full_text = ''
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + '. '
    return full_text

def get_text_xlsx(path):
    full_text = ''
    wb = load_workbook(path, data_only=True)
    for sheetname in wb.sheetnames:
        ws = wb[sheetname]
        all_rows = list(ws.rows)
        for row in all_rows:
            for r in row:
                if r.value is not None:
                    full_text += str(r.value) + '. '
    return full_text

def get_text_msg(path):
    msg = extract_msg.openMsg(path)
    full_text = msg.subject + '. ' + msg.body
    return full_text
